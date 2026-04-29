"""
Generate a .pptx deck summarizing the FirstAid-AI eval work.
Open the output in Google Slides or PowerPoint.

Usage: python3 llm/make_slides.py
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- colors ----------
BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)     # teal
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)    # coral/red
ACCENT3 = RGBColor(0xFF, 0xD9, 0x3D)    # yellow
DARK_CARD = RGBColor(0x25, 0x25, 0x3A)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_slide(slide, items, start_top=1.8, color=LIGHT_GRAY, size=16):
    for i, item in enumerate(items):
        add_text(slide, 0.8, start_top + i * 0.55, 8.4, 0.5,
                 f"  {item}", size=size, color=color)


def add_card(slide, left, top, width, height, color=DARK_CARD):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)  # rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Load results if available
    results_path = Path(__file__).parent / "results.json"
    results = None
    if results_path.exists():
        results = json.loads(results_path.read_text())

    # ===== SLIDE 1: Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 1.2, 9, 1.2, "FirstAid-AI", size=44, color=ACCENT, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 2.2, 9, 0.8, "On-Device First-Aid Triage with Local LLMs",
             size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.2, 9, 0.5, "Eval Pipeline Progress Update",
             size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ===== SLIDE 2: What We Built =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "What We Built This Week", size=30, color=ACCENT, bold=True)
    items = [
        "Eval harness (run_eval.py) — fully local, no API keys needed",
        "20 eval cases across 4 severity tiers (T0-T3)",
        "Structured grading: required / preferred / prohibited behaviors",
        "Local judge model (Qwen2.5-14B) scores each output",
        "Multi-model comparison: 7 models benchmarked in one run",
        "Automated scoring: weighted score + Tier 3 miss rate + VLM quality breakdown",
    ]
    add_bullet_slide(slide, items, start_top=1.2)

    # ===== SLIDE 3: Pipeline Architecture =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "Pipeline Architecture", size=30, color=ACCENT, bold=True)

    steps = [
        ("1  Load Cases", "eval_cases.json\n20 cases, 4 tiers"),
        ("2  Run Model", "mlx_lm.generate\n7 models tested"),
        ("3  Judge", "Qwen2.5-14B\nreq/pref/prohib grading"),
        ("4  Aggregate", "results.json\nper-tier + weighted scores"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = 0.4 + i * 2.4
        add_card(slide, x, 1.3, 2.1, 2.2)
        add_text(slide, x + 0.1, 1.4, 1.9, 0.5, title, size=15, color=ACCENT, bold=True)
        add_text(slide, x + 0.1, 1.9, 1.9, 1.4, desc, size=12, color=LIGHT_GRAY)
        if i < len(steps) - 1:
            add_text(slide, x + 2.1, 2.0, 0.3, 0.5, ">", size=24, color=ACCENT3, bold=True)

    add_text(slide, 0.5, 3.8, 9, 0.5,
             "Everything runs on-device (Apple Silicon) — no cloud, no API costs",
             size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ===== SLIDE 4: Tier System =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "Severity Tier System", size=30, color=ACCENT, bold=True)

    tiers = [
        ("Tier 0 — Self-Care", "Paper cuts, bruises, runny nose", "Weight: 1x", WHITE),
        ("Tier 1 — Moderate", "Sprains, minor burns, bee stings", "Weight: 2x", ACCENT3),
        ("Tier 2 — ER Required", "Fractures, deep cuts, head injuries", "Weight: 3x", RGBColor(0xFF, 0xA5, 0x00)),
        ("Tier 3 — Life-Threatening", "Heart attack, stroke, anaphylaxis", "Weight: 5x", ACCENT2),
    ]
    for i, (title, examples, weight, color) in enumerate(tiers):
        y = 1.1 + i * 1.0
        add_card(slide, 0.5, y, 9, 0.85)
        add_text(slide, 0.7, y + 0.05, 4, 0.4, title, size=16, color=color, bold=True)
        add_text(slide, 0.7, y + 0.4, 5.5, 0.4, examples, size=13, color=LIGHT_GRAY)
        add_text(slide, 8.0, y + 0.15, 1.3, 0.4, weight, size=14, color=color, bold=True,
                 align=PP_ALIGN.RIGHT)

    # ===== SLIDE 5: Eval Design Choices =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "Key Eval Design Choices", size=30, color=ACCENT, bold=True)
    items = [
        "Tier 3 users MINIMIZE symptoms — model must catch what user downplays",
        "VLM inputs vary in quality: 60% good, 25% vague, 10% contradictory, 5% useless",
        "Structured grading: required (must have) / preferred (bonus) / prohibited (auto-fail)",
        "Judge is 14B (4.7x larger than test models) — strong enough to grade 3B-8B",
        "Demographic tags: standard adult, pediatric, elderly",
        "Contradictory VLM cases test if model trusts user over bad image analysis",
    ]
    add_bullet_slide(slide, items, start_top=1.2)

    # ===== SLIDE 6: Models Under Test =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "Models Under Test", size=30, color=ACCENT, bold=True)

    models = [
        ("Llama 3.2 3B Instruct", "3B", "Meta — strong baseline"),
        ("Qwen3 4B", "4B", "Alibaba — good reasoning"),
        ("Phi-4 Mini", "3.8B", "Microsoft — compact performer"),
        ("Gemma 3 4B", "4B", "Google — instruction-tuned"),
        ("Ministral 3B", "3B", "Mistral — efficient"),
        ("SmolLM2 1.7B", "1.7B", "HuggingFace — floor test"),
        ("Qwen3 8B", "8B", "Alibaba — ceiling test"),
    ]
    # Header
    add_text(slide, 0.7, 1.1, 4, 0.4, "Model", size=14, color=ACCENT, bold=True)
    add_text(slide, 5.5, 1.1, 1.2, 0.4, "Size", size=14, color=ACCENT, bold=True)
    add_text(slide, 6.7, 1.1, 3, 0.4, "Notes", size=14, color=ACCENT, bold=True)

    for i, (name, size, notes) in enumerate(models):
        y = 1.55 + i * 0.48
        add_text(slide, 0.7, y, 4, 0.4, name, size=14, color=WHITE)
        add_text(slide, 5.5, y, 1.2, 0.4, size, size=14, color=ACCENT3)
        add_text(slide, 6.7, y, 3, 0.4, notes, size=12, color=LIGHT_GRAY)

    # ===== SLIDE 7: Llama Results (from actual data) =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "Llama 3.2 3B — Results (Baseline)", size=30, color=ACCENT, bold=True)

    if results and "summary" in results:
        summary = results["summary"]
        per_tier = summary.get("per_tier", {})

        # Tier results
        for i, t in enumerate(["0", "1", "2", "3"]):
            stats = per_tier.get(t, {"pass": 0, "fail": 0, "rate": 0})
            y = 1.2 + i * 0.65
            p, f = stats["pass"], stats["fail"]
            rate = stats["rate"]
            color = ACCENT if rate >= 0.8 else ACCENT3 if rate >= 0.6 else ACCENT2
            add_card(slide, 0.5, y, 5.5, 0.55)
            add_text(slide, 0.7, y + 0.05, 2, 0.4, f"Tier {t}", size=16, color=WHITE, bold=True)
            add_text(slide, 3.0, y + 0.05, 1.5, 0.4, f"{p}/{p+f} pass", size=15, color=color)
            add_text(slide, 4.5, y + 0.05, 1.2, 0.4, f"({rate:.0%})", size=15, color=color, bold=True)

        # Summary stats
        add_card(slide, 6.5, 1.2, 3, 1.1)
        add_text(slide, 6.7, 1.3, 2.6, 0.4, "Weighted Score", size=14, color=LIGHT_GRAY)
        ws = summary.get("weighted_score", 0)
        ws_color = ACCENT if ws >= 0.8 else ACCENT3 if ws >= 0.6 else ACCENT2
        add_text(slide, 6.7, 1.65, 2.6, 0.5, f"{ws:.0%}", size=32, color=ws_color, bold=True)

        add_card(slide, 6.5, 2.5, 3, 1.1)
        add_text(slide, 6.7, 2.6, 2.6, 0.4, "Tier 3 Miss Rate", size=14, color=LIGHT_GRAY)
        t3m = summary.get("tier_3_miss_rate", 0)
        t3_color = ACCENT if t3m <= 0.1 else ACCENT3 if t3m <= 0.3 else ACCENT2
        add_text(slide, 6.7, 2.95, 2.6, 0.5, f"{t3m:.0%}", size=32, color=t3_color, bold=True)

        # VLM quality
        vlm = summary.get("by_vlm_quality", {})
        add_text(slide, 0.5, 4.0, 9, 0.4, "Pass rate by VLM quality:", size=14, color=LIGHT_GRAY)
        vlm_parts = []
        for q in ["good", "vague", "contradictory", "useless"]:
            if q in vlm:
                vlm_parts.append(f"{q}: {vlm[q]['rate']:.0%}")
        add_text(slide, 0.5, 4.35, 9, 0.4, "    ".join(vlm_parts), size=14, color=ACCENT3)

    # ===== SLIDE 8: Key Findings =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "Key Findings So Far", size=30, color=ACCENT, bold=True)
    items = [
        "Model says 'CALL 911' for everything — even paper cuts (over-escalation)",
        "Tier 3: calls 911 but misses specifics (EpiPen, secondary drowning, stroke ID)",
        "Useless VLM input → model refused to give advice entirely (0% pass)",
        "Contradictory VLM input → 67% pass — model sometimes trusts bad image over user",
        "Judge passes cases with missing required behaviors — needs stricter calibration",
        "System prompt needs work: structured output format helps but doesn't fix triage logic",
    ]
    add_bullet_slide(slide, items, start_top=1.2)

    # ===== SLIDE 9: System Prompt Evolution =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "System Prompt — Current Version", size=30, color=ACCENT, bold=True)

    add_card(slide, 0.5, 1.1, 9, 3.5)
    prompt_text = (
        "Structured output format:\n"
        "  SEVERITY: LOW | MODERATE | EMERGENCY\n"
        "  ACTIONS: 1. / 2. / 3.\n"
        "  WATCH FOR: [worsening signs]\n\n"
        "Key rules:\n"
        "  - EMERGENCY → action 1 MUST be 'Call 911 now'\n"
        "  - Immediate actions only, no long-term care\n"
        "  - No greetings, apologies, or conversational language\n"
        "  - Ambiguous = EMERGENCY (err toward caution)\n"
        "  - Actions read aloud by TTS\n\n"
        "Includes 4 few-shot examples (LOW, MODERATE, 2x EMERGENCY)"
    )
    add_text(slide, 0.7, 1.2, 8.6, 3.3, prompt_text, size=13, color=LIGHT_GRAY)

    # ===== SLIDE 10: Next Steps =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_text(slide, 0.5, 0.3, 9, 0.7, "Next Steps", size=30, color=ACCENT, bold=True)

    items = [
        "Run full 7-model comparison and pick top 2-3 candidates",
        "Fix over-escalation: model should NOT call 911 for paper cuts",
        "Tighten judge: currently passes cases with missing required behaviors",
        "Scale eval to 200 cases for statistical significance",
        "Test with real VLM output (not simulated vlm_input strings)",
        "Begin MLX Swift integration for on-device iOS prototype",
    ]
    add_bullet_slide(slide, items, start_top=1.2, size=17)

    # ---------- Save ----------
    out_path = Path(__file__).parent / "FirstAid-AI_Eval_Update.pptx"
    prs.save(str(out_path))
    print(f"Slides saved to {out_path}")
    print("Open in Google Slides: upload to Google Drive > Open with Google Slides")


if __name__ == "__main__":
    make_deck()
